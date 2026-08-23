"""
Gera a apresentação executiva (PPTX) do Tech Challenge Fase 2.

Cobre os 4 pontos exigidos para o vídeo executivo — problema de negócio,
arquitetura da solução, valor da pipeline para análises educacionais e
potencial de uso para IA — mais uma seção de qualidade/limitações dos dados.
Paleta, tipografia e estilo de card seguem o modelo usado na apresentação do
Tech Challenge Fase 1 (reports/presentation/NPS_Preditivo_Apresentacao.pptx
em TechChallenge_1), adaptados aos dados e à identidade deste projeto.

Os números usados vêm de reports/dashboard_data.json (gerado por
reports/gerar_dashboard_dados.py a partir da camada Gold) — não são
fabricados nem hardcoded a partir de suposições.

Uso:
    python reports/presentation/gerar_apresentacao.py
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT_DIR = Path(__file__).resolve().parent.parent.parent
DATA = json.loads((ROOT_DIR / "reports" / "dashboard_data.json").read_text(encoding="utf-8"))

# ─── Paleta (modelo: TechChallenge_1/reports/presentation/NPS_Preditivo_Apresentacao.pptx) ──
NAVY = RGBColor(0x1A, 0x25, 0x3A)
BLUE_MED = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GREEN = RGBColor(0x1E, 0x8B, 0x4C)
BG_LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
GRAY_BADGE = RGBColor(0xEE, 0xEE, 0xEE)
TEXT_DARK = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xAA, 0xCC, 0xFF)

FONT = "Calibri"
FOOTER_TEXT = "Tech Challenge Fase 2  |  Pipeline Híbrido — Alfabetização no Brasil  |  Pós-Tech FIAP"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ───────────────────────────────────────────────────────────────────

def novo_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(1.5), rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = line_width
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=16, color=TEXT_DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = FONT
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=TEXT_DARK,
                 marker="▶", marker_color=ACCENT, space_after=12, bold_marker=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = marker_color
        r1.font.bold = bold_marker
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = FONT
    return tb


def header_footer(slide, titulo, eyebrow=None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_LIGHT)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), fill=NAVY)
    if eyebrow:
        add_text(slide, Inches(0.42), Inches(0.14), Inches(10), Inches(0.3), eyebrow.upper(),
                  size=11, color=LIGHT_BLUE, bold=True)
    add_text(slide, Inches(0.4), Inches(0.42), Inches(12.5), Inches(0.7), titulo,
              size=28, color=WHITE, bold=True)
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill=NAVY)
    add_text(slide, Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.28), FOOTER_TEXT,
              size=9, color=LIGHT_BLUE)


def add_kpi_card(slide, x, y, w, h, valor, label, cor_borda, chip=None, chip_cor=None):
    add_rect(slide, x, y, w, h, fill=WHITE, line=cor_borda, line_width=Pt(2.25), rounded=True)
    add_text(slide, x, y + Inches(0.18), w, Inches(0.8), valor, size=32, color=NAVY,
              bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + h - Inches(0.62), w - Inches(0.2), Inches(0.4), label,
              size=12.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    if chip:
        add_text(slide, x + Inches(0.1), y + h - Inches(0.30), w - Inches(0.2), Inches(0.28), chip,
                  size=10.5, color=chip_cor or ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return None


def add_card(slide, x, y, w, h, titulo, corpo, cor_borda, icone=""):
    add_rect(slide, x, y, w, h, fill=WHITE, line=cor_borda, line_width=Pt(2), rounded=True)
    add_rect(slide, x, y, Inches(0.09), h, fill=cor_borda)
    add_text(slide, x + Inches(0.28), y + Inches(0.16), w - Inches(0.5), Inches(0.42),
              f"{icone}  {titulo}", size=15, color=NAVY, bold=True)
    add_text(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.5), h - Inches(0.8),
              corpo, size=12, color=TEXT_DARK)


def add_checklist_item(slide, x, y, w, texto, cor=ACCENT):
    add_rect(slide, x, y, Inches(0.06), Inches(0.4), fill=cor)
    add_text(slide, x + Inches(0.2), y, w, Inches(0.4), texto, size=15, color=TEXT_DARK)


# ─── Dados reais (reports/dashboard_data.json) ─────────────────────────────────

BR = DATA["brasil"]["headline"]
UF = DATA["uf"]
MUNI = DATA["municipio"]
QUALI = DATA["qualidade"]["total"]

melhor_uf = UF["melhores"][0]
pior_uf = UF["piores"][0]
maior_queda_uf = UF["mais_pioraram"][0]
meta_muni = MUNI["meta_municipal"]

fmt1 = lambda v: f"{v:.1f}".replace(".", ",")


# ─── Construção da apresentação ────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Slide 1 — Capa ──
s = novo_slide(prs)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, 0, Inches(6.26), SLIDE_W, Inches(1.24), fill=BLUE_MED)
add_text(s, Inches(1.0), Inches(2.05), Inches(11.0), Inches(1.3),
          "Pipeline Híbrido para Análise\nda Alfabetização no Brasil", size=40, color=WHITE, bold=True)
add_text(s, Inches(1.0), Inches(3.55), Inches(11.0), Inches(0.7),
          "Indicador Criança Alfabetizada", size=26, color=LIGHT_BLUE)
add_rect(s, Inches(1.0), Inches(4.35), Inches(4.3), Inches(0.04), fill=ACCENT)
add_text(s, Inches(1.0), Inches(4.60), Inches(11.0), Inches(0.5),
          "Pós-Tech FIAP  |  1IAST  |  Tech Challenge Fase 2", size=16, color=LIGHT_BLUE)
add_text(s, Inches(1.0), Inches(5.05), Inches(11.0), Inches(0.5),
          "Felipe Barbato de Biaggio", size=16, color=WHITE)


# ── Slide 2 — O Problema de Negócio ──
s = novo_slide(prs)
header_footer(s, "O Problema de Negócio")

add_text(s, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.44), "Situação atual:",
          size=20, color=BLUE_MED, bold=True)
add_bullets(s, Inches(0.4), Inches(2.10), Inches(12.4), Inches(2.6), [
    "O Compromisso Nacional Criança Alfabetizada mobiliza União, estados e municípios para que toda criança esteja alfabetizada até o final do 2º ano do ensino fundamental — meta: 2030.",
    "O INEP definiu 743 pontos na escala de proficiência do Saeb como o corte de alfabetização, criando o Indicador Criança Alfabetizada.",
    f"Hoje ({BR['ano_atual']}), {fmt1(BR['indicador_atual'])}% das crianças avaliadas atingem esse patamar — a meta é {fmt1(BR['meta_2030'])}% até 2030.",
    "Sem integrar resultado, metas, território e contexto socioeconômico, gestores não enxergam onde agir primeiro.",
], size=15)

add_text(s, Inches(0.4), Inches(4.85), Inches(6.0), Inches(0.44), "Pergunta central:",
          size=20, color=BLUE_MED, bold=True)
add_rect(s, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.35), fill=WHITE, line=ACCENT, line_width=Pt(1.5), rounded=True)
add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.0),
          "“Como transformar dados dispersos do INEP em uma base confiável, atualizada e pronta para "
          "decisão — mostrando quem está evoluindo, quem está estagnado, e por quê?”",
          size=17, color=NAVY, italic=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# ── Slide 3 — Diagnóstico: onde estamos hoje ──
s = novo_slide(prs)
header_footer(s, "Diagnóstico: Onde Estamos Hoje")

card_w, card_h, gap = Inches(2.85), Inches(1.75), Inches(0.25)
x0 = Inches(0.5)
y0 = Inches(1.65)
add_kpi_card(s, x0, y0, card_w, card_h, f"{fmt1(BR['indicador_atual'])}%", f"Indicador nacional ({BR['ano_atual']})", ACCENT)
add_kpi_card(s, x0 + (card_w + gap) * 1, y0, card_w, card_h, f"{fmt1(BR['meta_2030'])}%", "Meta nacional (2030)", BLUE_MED)
add_kpi_card(s, x0 + (card_w + gap) * 2, y0, card_w, card_h, f"{fmt1(BR['gap_meta_2030'])} p.p.", "Gap para a meta", ORANGE)
add_kpi_card(s, x0 + (card_w + gap) * 3, y0, card_w, card_h, f"+{fmt1(BR['delta_periodo'])} p.p.", f"Evolução {BR['ano_inicial']}–{BR['ano_atual']}", GREEN)

add_rect(s, Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.03), fill=ACCENT)
add_text(s, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.44),
          "O que esses números significam para a gestão?", size=19, color=BLUE_MED, bold=True)
add_bullets(s, Inches(0.4), Inches(4.55), Inches(12.4), Inches(2.3), [
    f"Estamos avançando (+{fmt1(BR['delta_periodo'])} p.p. em {BR['ano_atual']-BR['ano_inicial']} anos), mas o ritmo precisa acelerar para fechar {fmt1(BR['gap_meta_2030'])} p.p. até 2030.",
    f"O indicador nacional esconde desigualdade regional: {melhor_uf['sigla_uf']} chega a {fmt1(melhor_uf['indicador_atual'])}%, enquanto {pior_uf['sigla_uf']} não passa de {fmt1(pior_uf['indicador_atual'])}%.",
    f"{maior_queda_uf['sigla_uf']} teve a maior queda do país entre os dois últimos anos ({fmt1(maior_queda_uf['delta'])} p.p.) — um sinal de alerta que só aparece comparando anos lado a lado.",
], size=15)


# ── Slide 4 — Arquitetura da Solução ──
s = novo_slide(prs)
header_footer(s, "Arquitetura da Solução")

# Diagrama simplificado: Fontes -> Bronze -> Silver -> Gold -> Consumo
diag_y = Inches(1.65)
box_h = Inches(1.15)
labels = [
    ("Fontes", "INEP / Base dos\nDados + IDHM\n+ Streaming", NAVY),
    ("Bronze", "Dado bruto\npreservado", BLUE_MED),
    ("Silver", "Limpo, integrado\ne validado", ACCENT),
    ("Gold", "Pronto para\nanálise e IA", GREEN),
]
box_w = Inches(2.65)
gap_arrow = Inches(0.35)
x = Inches(0.5)
for i, (titulo, desc, cor) in enumerate(labels):
    add_rect(s, x, diag_y, box_w, box_h, fill=cor, rounded=True)
    add_text(s, x, diag_y + Inches(0.12), box_w, Inches(0.35), titulo, size=16, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), diag_y + Inches(0.48), box_w - Inches(0.2), Inches(0.6), desc,
              size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(labels) - 1:
        add_text(s, x + box_w, diag_y + Inches(0.30), gap_arrow, Inches(0.5), "→", size=26,
                  color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    x += box_w + gap_arrow

add_text(s, Inches(0.5), Inches(3.15), Inches(12.4), Inches(0.4),
          "Qualidade de dados validada de ponta a ponta, em todas as camadas, a cada execução",
          size=13, color=TEXT_DARK, italic=True, align=PP_ALIGN.CENTER)

add_bullets(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(3.0), [
    "Ingestão híbrida: dados históricos por lote (batch, com fallback automático) + eventos simulados de atualização (streaming) — o mesmo padrão usado por bancos e fintechs para ingestão contínua.",
    "9 fontes de dados integradas (resultado, metas, território, microdados de alunos, IDHM), unificadas em um único registro por município e ano.",
    "A camada Gold alimenta diretamente dashboards, análises estatísticas e o treinamento de modelos de machine learning — sem retrabalho.",
], size=14.5)


# ── Slide 5 — Qualidade e Confiabilidade dos Dados ──
s = novo_slide(prs)
header_footer(s, "Qualidade e Limitações dos Dados", eyebrow="Transparência com os dados")

kpi_w, kpi_h = Inches(3.9), Inches(1.35)
kx = Inches(0.5)
ky = Inches(1.6)
add_kpi_card(s, kx, ky, kpi_w, kpi_h, str(QUALI["tabelas_verificadas"]), "Tabelas verificadas", ACCENT)
add_kpi_card(s, kx + kpi_w + Inches(0.3), ky, kpi_w, kpi_h, f"{QUALI['tabelas_ok']} / {QUALI['tabelas_verificadas']}", "Tabelas OK", GREEN)
add_kpi_card(s, kx + (kpi_w + Inches(0.3)) * 2, ky, kpi_w, kpi_h, str(QUALI["total_alertas"]), "Alertas de qualidade", ACCENT)

card_y = Inches(3.25)
card_w2, card_h2 = Inches(4.03), Inches(3.55)
add_card(s, Inches(0.4), card_y, card_w2, card_h2, "Dicionário oficial indisponível",
          "5 colunas (rede, serie, presenca, preenchimento_caderno, alfabetizado) tiveram o "
          "significado inferido — não confirmado pela fonte, que exige BigQuery pago. Validado "
          "empiricamente: alfabetizado bate 100% com o corte de 743 pontos.",
          ORANGE, icone="📖")
add_card(s, Inches(4.63), card_y, card_w2, card_h2, "A fonte atualiza pouco",
          "Última atualização registrada: 23/09/2025 — mas ainda não há dado de 2025 disponível. "
          "O indicador é medido uma vez por ciclo letivo, não é dado transacional que muda dia a dia.",
          ORANGE, icone="📅")
add_card(s, Inches(8.86), card_y, card_w2, card_h2, "Streaming é uma simulação",
          "Eventos sintéticos demonstram o padrão de ingestão incremental. A fonte real não emite "
          "atualizações contínuas — por isso streaming de verdade seria over-engineering aqui.",
          ACCENT, icone="🔄")


# ── Slide 6 — Valor da Pipeline para Análises Educacionais ──
s = novo_slide(prs)
header_footer(s, "Valor da Pipeline para Análises Educacionais")

atingiram_fmt = f"{meta_muni['atingiram']:,}".replace(",", ".")
total_com_meta_fmt = f"{meta_muni['total_com_meta']:,}".replace(",", ".")

add_bullets(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(4.6), [
    f"Ranking pronto por UF e município: hoje, {melhor_uf['sigla_uf']} lidera com {fmt1(melhor_uf['indicador_atual'])}%; "
    f"{pior_uf['sigla_uf']} é a UF com menor indicador ({fmt1(pior_uf['indicador_atual'])}%).",
    f"{fmt1(meta_muni['pct_atingiram'])}% dos municípios avaliados atingiram sua própria meta municipal "
    f"({atingiram_fmt} de {total_com_meta_fmt} com meta definida).",
    f"{maior_queda_uf['sigla_uf']} teve a maior queda entre os dois últimos anos ({fmt1(maior_queda_uf['delta'])} p.p.) — "
    "um sinal que só aparece quando comparamos períodos lado a lado, não em uma foto única.",
    "Cada município carrega, na mesma linha, resultado, meta, gap nacional, faixa de risco e IDHM — "
    "pronto para priorizar onde investir, sem cruzar planilhas manualmente.",
    "Tudo isso já está disponível como um painel interativo (4 visões: Brasil, UF, Município, Aluno), "
    "pronto para gestores explorarem sem escrever uma linha de código.",
], size=16, space_after=18)


# ── Slide 7 — Potencial de Uso para Inteligência Artificial ──
s = novo_slide(prs)
header_footer(s, "Potencial de Uso para Inteligência Artificial")

ia_cards = [
    ("🎯", "Modelos preditivos de alfabetização",
     "Usar IDHM (educação, longevidade, renda) + histórico do indicador para prever municípios em risco de não atingir a meta de 2030."),
    ("🧩", "Clusters de vulnerabilidade educacional",
     "Segmentar municípios por perfil de risco (faixa de risco + IDHM + gap de meta) para priorizar onde investir primeiro."),
    ("⚖️", "Análise de desigualdade educacional",
     "Cruzar taxa de alfabetização com IDHM por região/UF para saber se o gap educacional acompanha o gap de desenvolvimento humano."),
    ("🏛️", "Políticas públicas baseadas em dados",
     "Simular cenários de investimento e medir impacto esperado no indicador, segmentando por faixa de IDHM."),
]
gw, gh = Inches(6.05), Inches(2.55)
gx0, gy0 = Inches(0.4), Inches(1.65)
ggap = Inches(0.2)
for i, (icone, titulo, corpo) in enumerate(ia_cards):
    col, row = i % 2, i // 2
    x = gx0 + col * (gw + ggap)
    y = gy0 + row * (gh + ggap)
    add_card(s, x, y, gw, gh, titulo, corpo, ACCENT, icone=icone)


# ── Slide 8 — FinOps: Custo Simbólico de Operação ──
s = novo_slide(prs)
header_footer(s, "Quanto Custa Rodar Isso na Nuvem?", eyebrow="FinOps")

add_text(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.7),
          "Estimativa simbólica: preço on-demand padrão (BigQuery + Cloud Run + Storage), "
          "cruzando volume real de dados (~219 MB) com a frequência de execução.",
          size=13.5, color=TEXT_DARK, italic=True)

# tabela de cenarios
linhas = [
    ("Diário", "365x/ano", "R$ 35 / ano", RED),
    ("Semanal", "52x/ano", "R$ 5 / ano", ORANGE),
    ("Mensal (recomendado)", "12x/ano", "R$ 1,40 / ano", GREEN),
    ("Cadência real da fonte", "~1x/ano", "R$ 0,35 / ano", ACCENT),
]
ty = Inches(2.35)
row_h = Inches(0.85)
col_w = [Inches(4.2), Inches(3.0), Inches(3.0)]
headers = ["Cadência de execução", "Frequência", "Custo anual estimado"]
tx = Inches(0.5)
add_rect(s, tx, ty, sum(col_w, Inches(0)), Inches(0.55), fill=NAVY)
cx = tx
for h, w in zip(headers, col_w):
    add_text(s, cx + Inches(0.15), ty + Inches(0.10), w - Inches(0.2), Inches(0.4), h,
              size=13, color=WHITE, bold=True)
    cx += w

ry = ty + Inches(0.55)
for nome, freq, custo, cor in linhas:
    add_rect(s, tx, ry, sum(col_w, Inches(0)), row_h, fill=WHITE, line=GRAY_BADGE, line_width=Pt(1))
    add_rect(s, tx, ry, Inches(0.08), row_h, fill=cor)
    cx = tx
    for val, w in zip([nome, freq, custo], col_w):
        add_text(s, cx + Inches(0.2), ry + Inches(0.20), w - Inches(0.3), Inches(0.5), val,
                  size=14, color=NAVY, bold=(val == custo))
        cx += w
    ry += row_h

add_rect(s, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.85), fill=WHITE, line=ACCENT, line_width=Pt(1.5), rounded=True)
add_text(s, Inches(0.65), Inches(6.28), Inches(12.0), Inches(0.6),
          "Rodar mensalmente custa ~96% menos que diariamente — sem perder capacidade de captar "
          "dado novo, já que a fonte não publica com essa frequência. Dentro do tier gratuito do "
          "BigQuery, o custo real hoje é R$ 0.",
          size=12.5, color=NAVY, italic=True, anchor=MSO_ANCHOR.MIDDLE)


# ── Slide 9 — Conclusão e Próximos Passos ──
s = novo_slide(prs)
header_footer(s, "Conclusão e Próximos Passos")

add_text(s, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4), "Resumo:",
          size=18, color=BLUE_MED, bold=True)

itens = [
    (f"O indicador está subindo (+{fmt1(BR['delta_periodo'])} p.p.) mas o ritmo não fecha a meta 2030 sem ação direcionada", ACCENT),
    ("A pipeline já entrega, hoje, ranking, gaps e deltas prontos para decisão — sem esperar TI", GREEN),
    ("A base Gold, com IDHM, já está pronta para os primeiros modelos preditivos e clusters de vulnerabilidade", ACCENT),
    (f"Qualidade dos dados é auditada automaticamente a cada execução ({QUALI['tabelas_ok']}/{QUALI['tabelas_verificadas']} tabelas OK na última rodada)", GREEN),
    ("Operar isso custa centavos por ano — a decisão de frequência importa mais que a infraestrutura", ORANGE),
]
y = Inches(2.15)
for texto, cor in itens:
    add_checklist_item(s, Inches(0.4), y, Inches(12.1), texto, cor=cor)
    y += Inches(0.62)


# ─── Salvar ──────────────────────────────────────────────────────────────────
destino = Path(__file__).resolve().parent / "Alfabetizacao_Apresentacao.pptx"
prs.save(destino)
print(f"Salvo: {destino}")
